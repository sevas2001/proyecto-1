"""
ModerApp — Generador de presentación PowerPoint profesional
Dependencia: pip install python-pptx
Uso:         python generate_pptx.py
Salida:      ModerApp_NTT.pptx  (en la raíz del proyecto)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paleta de colores ──────────────────────────────────────────────────────────
NAVY       = RGBColor(0x00, 0x2B, 0x5C)   # azul oscuro NTT
BLUE       = RGBColor(0x00, 0x5B, 0xAA)   # azul medio
RED_ACCENT = RGBColor(0xE8, 0x30, 0x2C)   # rojo acento NTT
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF0, 0xF4, 0xF9)   # fondo claro diapositivas internas
GRAY_TEXT  = RGBColor(0x55, 0x55, 0x66)
YELLOW_HL  = RGBColor(0xFF, 0xC8, 0x00)   # amarillo highlight

# ── Dimensiones (16:9) ─────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

DIAGRAM_PATH = os.path.join(os.path.dirname(__file__), "diagrama.png")

# ──────────────────────────────────────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────────────────────────────────────

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # Layout en blanco


def rect(slide, left, top, width, height, fill_color, alpha=None):
    """Rectángulo sólido sin borde."""
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp


def textbox(slide, left, top, width, height, text,
            size=18, bold=False, color=WHITE,
            align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = text
    run.font.size       = Pt(size)
    run.font.bold       = bold
    run.font.italic     = italic
    run.font.color.rgb  = color
    return txb


def bullet_frame(slide, left, top, width, height,
                 items, size=16, color=WHITE,
                 bullet_char="▸", line_spacing=1.2):
    """Cuadro de texto con lista de bullets."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text           = f"{bullet_char}  {item}"
        run.font.size      = Pt(size)
        run.font.color.rgb = color
        # espaciado
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", str(int(line_spacing * 100000)))
    return txb


def badge(slide, left, top, width, height, label, bg=BLUE, fg=WHITE, size=14):
    """Pastilla de color con texto centrado."""
    shp = rect(slide, left, top, width, height, bg)
    # borde redondeado via XML
    spPr = shp.fill._xPr.getparent().getparent()
    prstGeom = spPr.find(qn("a:prstGeom"))
    if prstGeom is not None:
        prstGeom.set("prst", "roundRect")
    txb = textbox(slide, left, top, width, height,
                  label, size=size, bold=True, color=fg,
                  align=PP_ALIGN.CENTER)
    return shp


# ──────────────────────────────────────────────────────────────────────────────
# DIAPOSITIVA 1 — PORTADA
# ──────────────────────────────────────────────────────────────────────────────

def slide_portada(prs):
    sl = blank_slide(prs)

    # Fondo completo navy
    rect(sl, 0, 0, W, H, NAVY)

    # Franja lateral izquierda azul
    rect(sl, 0, 0, Inches(0.45), H, BLUE)

    # Línea decorativa roja (horizontal, parte superior)
    rect(sl, Inches(0.7), Inches(1.6), Inches(5.5), Inches(0.07), RED_ACCENT)

    # Título principal
    textbox(sl, Inches(0.7), Inches(1.8), Inches(8), Inches(1.6),
            "ModerApp",
            size=64, bold=True, color=WHITE)

    # Subtítulo
    textbox(sl, Inches(0.7), Inches(3.3), Inches(9), Inches(0.8),
            "Sistema de Moderación de Contenido con Inteligencia Artificial",
            size=24, color=RGBColor(0xAA, 0xCC, 0xEE))

    # Descripción breve
    textbox(sl, Inches(0.7), Inches(4.15), Inches(8.5), Inches(0.6),
            "Pipeline multi-agente · Azure GPT-4o · LangGraph · Streamlit",
            size=16, italic=True, color=RGBColor(0x88, 0xAA, 0xCC))

    # Línea separadora
    rect(sl, Inches(0.7), Inches(4.85), Inches(9), Inches(0.04),
         RGBColor(0x33, 0x55, 0x88))

    # Empresa / equipo
    textbox(sl, Inches(0.7), Inches(5.0), Inches(7), Inches(0.5),
            "NTT Data  ·  2025",
            size=14, color=RGBColor(0x77, 0x99, 0xBB))

    # Etiqueta "IA" decorativa (esquina inferior derecha)
    rect(sl, Inches(10.8), Inches(5.8), Inches(2.0), Inches(1.1), BLUE)
    textbox(sl, Inches(10.8), Inches(5.8), Inches(2.0), Inches(1.1),
            "AI\nMODERATION",
            size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# DIAPOSITIVA 2 — PROBLEMA & SOLUCIÓN
# ──────────────────────────────────────────────────────────────────────────────

def slide_problema(prs):
    sl = blank_slide(prs)

    # Fondo claro
    rect(sl, 0, 0, W, H, LIGHT_BG)

    # Cabecera navy
    rect(sl, 0, 0, W, Inches(1.2), NAVY)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(10), Inches(0.7),
            "El Problema & La Solución",
            size=30, bold=True, color=WHITE)

    # Número de diapositiva
    textbox(sl, Inches(12.3), Inches(0.35), Inches(0.8), Inches(0.5),
            "01", size=22, bold=True, color=YELLOW_HL, align=PP_ALIGN.RIGHT)

    # ── Panel izquierdo: Problema ─────────────────────────────────────────────
    rect(sl, Inches(0.4), Inches(1.5), Inches(5.8), Inches(4.8),
         RGBColor(0xFF, 0xEB, 0xEB))

    textbox(sl, Inches(0.55), Inches(1.65), Inches(5.5), Inches(0.6),
            "El Problema",
            size=20, bold=True, color=RED_ACCENT)

    bullet_frame(sl,
        Inches(0.55), Inches(2.3), Inches(5.5), Inches(3.5),
        items=[
            "La moderación manual es lenta,\ncostosa e inconsistente",
            "Escala imposible ante millones\nde publicaciones diarias",
            "Errores humanos generan riesgos\nlegales y reputacionales",
            "Contenido dañino (odio, spam,\nacoso) se propaga sin control",
        ],
        size=15, color=RGBColor(0x44, 0x11, 0x11), line_spacing=1.4)

    # ── Panel derecho: Solución ───────────────────────────────────────────────
    rect(sl, Inches(6.6), Inches(1.5), Inches(6.3), Inches(4.8),
         RGBColor(0xE6, 0xF0, 0xFF))

    textbox(sl, Inches(6.75), Inches(1.65), Inches(6.0), Inches(0.6),
            "Nuestra Solución",
            size=20, bold=True, color=BLUE)

    bullet_frame(sl,
        Inches(6.75), Inches(2.3), Inches(6.0), Inches(3.5),
        items=[
            "Pipeline automático con 3 agentes IA\norquestados por LangGraph",
            "Clasificación multimodal texto + imagen\ncon Azure GPT-4o",
            "Decisión basada en política YAML\npersonalizable",
            "Escalado a revisión humana cuando\nel modelo no tiene confianza suficiente",
        ],
        size=15, color=RGBColor(0x00, 0x22, 0x55), line_spacing=1.4)

    # Flecha central (representada con texto)
    textbox(sl, Inches(5.8), Inches(3.2), Inches(0.8), Inches(1.2),
            "→", size=48, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Footer
    rect(sl, 0, Inches(6.9), W, Inches(0.6), NAVY)
    textbox(sl, Inches(0.5), Inches(6.95), Inches(8), Inches(0.45),
            "ModerApp  |  NTT Data", size=11, color=RGBColor(0x88, 0xAA, 0xCC))


# ──────────────────────────────────────────────────────────────────────────────
# DIAPOSITIVA 3 — ARQUITECTURA
# ──────────────────────────────────────────────────────────────────────────────

def slide_arquitectura(prs):
    sl = blank_slide(prs)

    rect(sl, 0, 0, W, H, LIGHT_BG)

    # Cabecera
    rect(sl, 0, 0, W, Inches(1.2), NAVY)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(10), Inches(0.7),
            "Arquitectura del Sistema",
            size=30, bold=True, color=WHITE)
    textbox(sl, Inches(12.3), Inches(0.35), Inches(0.8), Inches(0.5),
            "02", size=22, bold=True, color=YELLOW_HL, align=PP_ALIGN.RIGHT)

    # Diagrama (imagen)
    if os.path.exists(DIAGRAM_PATH):
        sl.shapes.add_picture(
            DIAGRAM_PATH,
            Inches(0.3), Inches(1.3),
            Inches(8.2), Inches(5.5)
        )
    else:
        textbox(sl, Inches(0.5), Inches(2.5), Inches(8), Inches(1.0),
                "[diagrama.png no encontrado]", size=14, color=GRAY_TEXT)

    # Panel lateral derecho — leyenda
    rect(sl, Inches(8.8), Inches(1.3), Inches(4.2), Inches(5.5),
         RGBColor(0xE6, 0xF0, 0xFF))

    textbox(sl, Inches(8.95), Inches(1.45), Inches(3.9), Inches(0.5),
            "Flujo de moderación",
            size=16, bold=True, color=NAVY)

    steps = [
        ("1", "Usuario crea post", BLUE),
        ("2", "Orchestrator\n(LangGraph)", NAVY),
        ("3", "Agent Classifier\n(texto + imagen)", RGBColor(0x5A, 0x2D, 0x82)),
        ("4", "Agent Decider\n(política YAML)", RGBColor(0x5A, 0x2D, 0x82)),
        ("5", "Agent Reviewer\n(decisión final)", RGBColor(0x5A, 0x2D, 0x82)),
        ("6", "APROBADO / RECHAZADO\n/ REVISIÓN HUMANA", RED_ACCENT),
    ]

    y = Inches(2.05)
    for num, label, col in steps:
        rect(sl, Inches(8.95), y, Inches(0.4), Inches(0.55), col)
        textbox(sl, Inches(8.95), y, Inches(0.4), Inches(0.55),
                num, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(sl, Inches(9.45), y, Inches(3.4), Inches(0.55),
                label, size=12, color=RGBColor(0x11, 0x22, 0x44))
        y += Inches(0.72)

    # Footer
    rect(sl, 0, Inches(6.9), W, Inches(0.6), NAVY)
    textbox(sl, Inches(0.5), Inches(6.95), Inches(8), Inches(0.45),
            "ModerApp  |  NTT Data", size=11, color=RGBColor(0x88, 0xAA, 0xCC))


# ──────────────────────────────────────────────────────────────────────────────
# DIAPOSITIVA 4 — STACK TECNOLÓGICO
# ──────────────────────────────────────────────────────────────────────────────

def slide_stack(prs):
    sl = blank_slide(prs)

    rect(sl, 0, 0, W, H, LIGHT_BG)

    # Cabecera
    rect(sl, 0, 0, W, Inches(1.2), NAVY)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(10), Inches(0.7),
            "Stack Tecnológico",
            size=30, bold=True, color=WHITE)
    textbox(sl, Inches(12.3), Inches(0.35), Inches(0.8), Inches(0.5),
            "03", size=22, bold=True, color=YELLOW_HL, align=PP_ALIGN.RIGHT)

    # 3 columnas
    cols = [
        {
            "title":  "Frontend / UI",
            "color":  RGBColor(0x00, 0x5B, 0xAA),
            "items":  [
                "Streamlit 1.x",
                "4 pestañas: Crear Post,\nFeed, Moderación, Auditoría",
                "Soporte multimodal\n(texto + imagen)",
                "Estado reactivo con\nst.session_state",
            ],
            "left": Inches(0.4),
        },
        {
            "title":  "IA & Orquestación",
            "color":  RGBColor(0x5A, 0x2D, 0x82),
            "items":  [
                "LangGraph (state machine)",
                "LangChain / LangChain-OpenAI",
                "Azure OpenAI GPT-4o\n(multimodal)",
                "3 agentes especializados\nclasifier · decider · reviewer",
                "Política configurable\nen YAML",
            ],
            "left": Inches(4.7),
        },
        {
            "title":  "Backend & Datos",
            "color":  RGBColor(0x00, 0x7A, 0x50),
            "items":  [
                "FastAPI + Uvicorn",
                "Pydantic v2 (validación)",
                "Persistencia JSON\n(escritura atómica)",
                "Auditoría JSONL\n(logs por agente)",
                "python-dotenv\n(.env con credenciales)",
            ],
            "left": Inches(9.0),
        },
    ]

    for col in cols:
        x = col["left"]
        # Header de columna
        rect(sl, x, Inches(1.4), Inches(3.9), Inches(0.65), col["color"])
        textbox(sl, x, Inches(1.4), Inches(3.9), Inches(0.65),
                col["title"], size=17, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
        # Cuerpo
        rect(sl, x, Inches(2.05), Inches(3.9), Inches(4.55),
             RGBColor(0xFF, 0xFF, 0xFF))
        bullet_frame(sl, x + Inches(0.15), Inches(2.15),
                     Inches(3.65), Inches(4.3),
                     col["items"], size=14,
                     color=RGBColor(0x11, 0x22, 0x44),
                     bullet_char="✓", line_spacing=1.35)

    # Footer
    rect(sl, 0, Inches(6.9), W, Inches(0.6), NAVY)
    textbox(sl, Inches(0.5), Inches(6.95), Inches(8), Inches(0.45),
            "ModerApp  |  NTT Data", size=11, color=RGBColor(0x88, 0xAA, 0xCC))


# ──────────────────────────────────────────────────────────────────────────────
# DIAPOSITIVA 5 — DEMO EN VIVO
# ──────────────────────────────────────────────────────────────────────────────

def slide_demo(prs):
    sl = blank_slide(prs)

    # Fondo navy para contraste
    rect(sl, 0, 0, W, H, NAVY)

    # Franja lateral
    rect(sl, 0, 0, Inches(0.45), H, BLUE)

    # Etiqueta grande "DEMO"
    textbox(sl, Inches(0.7), Inches(0.8), Inches(12), Inches(2.0),
            "Demo en Vivo",
            size=58, bold=True, color=WHITE)

    # Línea roja decorativa
    rect(sl, Inches(0.7), Inches(2.65), Inches(10), Inches(0.07), RED_ACCENT)

    # Qué vamos a mostrar
    textbox(sl, Inches(0.7), Inches(2.85), Inches(5.5), Inches(0.55),
            "¿Qué veremos?",
            size=20, bold=True, color=YELLOW_HL)

    pasos = [
        "Crear un post con texto e imagen desde la UI",
        "Pipeline de moderación: Classifier → Decider → Reviewer",
        "Resultado en tiempo real: APROBADO / RECHAZADO / REVISIÓN HUMANA",
        "Log de auditoría con trazabilidad completa por agente",
        "Caso edge: contenido con baja confianza → escalado humano",
    ]
    bullet_frame(sl, Inches(0.7), Inches(3.45), Inches(8.5), Inches(3.0),
                 pasos, size=16, color=RGBColor(0xCC, 0xDD, 0xFF),
                 bullet_char="→", line_spacing=1.4)

    # Panel derecho — URL local
    rect(sl, Inches(9.5), Inches(2.8), Inches(3.4), Inches(2.2),
         RGBColor(0x00, 0x3A, 0x80))
    textbox(sl, Inches(9.5), Inches(2.85), Inches(3.4), Inches(0.55),
            "App local",
            size=15, bold=True, color=YELLOW_HL, align=PP_ALIGN.CENTER)
    textbox(sl, Inches(9.5), Inches(3.4), Inches(3.4), Inches(0.5),
            "http://localhost:8501",
            size=13, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, Inches(9.5), Inches(3.9), Inches(3.4), Inches(0.9),
            "Streamlit · Python 3.11\nAzure GPT-4o activo",
            size=12, color=RGBColor(0x88, 0xCC, 0xFF), align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# DIAPOSITIVA 6 — CONCLUSIONES & PRÓXIMOS PASOS
# ──────────────────────────────────────────────────────────────────────────────

def slide_conclusiones(prs):
    sl = blank_slide(prs)

    rect(sl, 0, 0, W, H, LIGHT_BG)

    # Cabecera
    rect(sl, 0, 0, W, Inches(1.2), NAVY)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(10), Inches(0.7),
            "Conclusiones & Próximos Pasos",
            size=30, bold=True, color=WHITE)
    textbox(sl, Inches(12.3), Inches(0.35), Inches(0.8), Inches(0.5),
            "05", size=22, bold=True, color=YELLOW_HL, align=PP_ALIGN.RIGHT)

    # ── Columna izquierda: Conclusiones ──────────────────────────────────────
    textbox(sl, Inches(0.5), Inches(1.45), Inches(5.8), Inches(0.55),
            "Conclusiones",
            size=19, bold=True, color=NAVY)
    rect(sl, Inches(0.5), Inches(1.95), Inches(5.8), Inches(0.05), BLUE)

    conclusiones = [
        "Pipeline multi-agente funcional y extensible",
        "Integración real con Azure GPT-4o (texto e imagen)",
        "Política de moderación configurable sin tocar código",
        "Auditoría completa con trazabilidad por agente",
        "Arquitectura lista para producción (FastAPI + Streamlit)",
    ]
    bullet_frame(sl, Inches(0.5), Inches(2.1), Inches(5.8), Inches(3.5),
                 conclusiones, size=15,
                 color=RGBColor(0x11, 0x22, 0x44), line_spacing=1.45)

    # ── Columna derecha: Próximos pasos ──────────────────────────────────────
    textbox(sl, Inches(7.0), Inches(1.45), Inches(5.8), Inches(0.55),
            "Próximos Pasos",
            size=19, bold=True, color=NAVY)
    rect(sl, Inches(7.0), Inches(1.95), Inches(5.8), Inches(0.05), RED_ACCENT)

    proximos = [
        "Integrar base de datos (PostgreSQL / MongoDB)",
        "Dashboard de métricas en tiempo real",
        "Fine-tuning del modelo con datos propios",
        "Sistema de roles y autenticación de moderadores",
        "Despliegue en Azure Container Apps",
    ]
    bullet_frame(sl, Inches(7.0), Inches(2.1), Inches(5.8), Inches(3.5),
                 proximos, size=15,
                 color=RGBColor(0x11, 0x22, 0x44), line_spacing=1.45)

    # ── Cierre ────────────────────────────────────────────────────────────────
    rect(sl, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.75),
         RGBColor(0xE0, 0xEB, 0xFF))
    textbox(sl, Inches(0.6), Inches(5.9), Inches(12.0), Inches(0.6),
            "ModerApp demuestra que la IA puede automatizar la moderación de contenido "
            "de forma escalable, explicable y segura.",
            size=14, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Footer
    rect(sl, 0, Inches(6.9), W, Inches(0.6), NAVY)
    textbox(sl, Inches(0.5), Inches(6.95), Inches(8), Inches(0.45),
            "ModerApp  |  NTT Data", size=11, color=RGBColor(0x88, 0xAA, 0xCC))


# ──────────────────────────────────────────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Generando diapositivas...")
    slide_portada(prs);      print("  [1/6] Portada OK")
    slide_problema(prs);     print("  [2/6] Problema & Solucion OK")
    slide_arquitectura(prs); print("  [3/6] Arquitectura OK")
    slide_stack(prs);        print("  [4/6] Stack tecnologico OK")
    slide_demo(prs);         print("  [5/6] Demo en vivo OK")
    slide_conclusiones(prs); print("  [6/6] Conclusiones OK")

    out_path = os.path.join(os.path.dirname(__file__), "ModerApp_NTT.pptx")
    prs.save(out_path)
    print(f"\nPresentación guardada en:\n  {out_path}")


if __name__ == "__main__":
    main()
